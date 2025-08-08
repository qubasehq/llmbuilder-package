"""
Data loading utilities for LLMBuilder.

This module provides the DataLoader class for loading various file formats
including TXT, PDF, DOCX, PPTX, and Markdown files.
"""

import os
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Iterator
from dataclasses import dataclass

from ..utils import DataError, get_logger

logger = get_logger("data.loader")

# Optional imports for different file formats
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class DocumentMetadata:
    """Metadata for a loaded document."""
    filename: str
    file_path: str
    file_size: int
    file_type: str
    char_count: int
    word_count: int
    line_count: int


class DataLoader:
    """
    Loads and processes various document formats for LLM training.
    
    Supports TXT, PDF, DOCX, PPTX, and Markdown files with automatic
    format detection and text extraction.
    """
    
    SUPPORTED_EXTENSIONS = {'.txt', '.pdf', '.docx', '.pptx', '.md', '.markdown'}
    
    def __init__(self, 
                 min_length: int = 50,
                 max_length: Optional[int] = None,
                 clean_text: bool = True):
        """
        Initialize the data loader.
        
        Args:
            min_length: Minimum text length to keep
            max_length: Maximum text length (None for no limit)
            clean_text: Whether to clean extracted text
        """
        self.min_length = min_length
        self.max_length = max_length
        self.clean_text = clean_text
        
        # Check for optional dependencies
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check for optional dependencies and warn if missing."""
        missing = []
        if not HAS_PDF:
            missing.append("PyPDF2 (for PDF support)")
        if not HAS_DOCX:
            missing.append("python-docx (for DOCX support)")
        if not HAS_PPTX:
            missing.append("python-pptx (for PPTX support)")
        
        if missing:
            logger.warning(f"Missing optional dependencies: {', '.join(missing)}")
    
    def load_file(self, file_path: Union[str, Path]) -> Optional[str]:
        """
        Load text content from a single file.
        
        Args:
            file_path: Path to the file to load
            
        Returns:
            str: Extracted text content, or None if extraction failed
            
        Raises:
            DataError: If file format is unsupported or loading fails
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise DataError(f"File not found: {file_path}")
        
        if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            raise DataError(f"Unsupported file format: {file_path.suffix}")
        
        try:
            # Route to appropriate loader based on extension
            if file_path.suffix.lower() == '.txt':
                text = self._load_txt(file_path)
            elif file_path.suffix.lower() == '.pdf':
                text = self._load_pdf(file_path)
            elif file_path.suffix.lower() == '.docx':
                text = self._load_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                text = self._load_pptx(file_path)
            elif file_path.suffix.lower() in ['.md', '.markdown']:
                text = self._load_markdown(file_path)
            else:
                raise DataError(f"Unsupported file format: {file_path.suffix}")
            
            if text and self.clean_text:
                text = self._clean_text(text)
            
            # Check length constraints
            if text and len(text) < self.min_length:
                logger.debug(f"Skipping {file_path.name}: too short ({len(text)} chars)")
                return None
            
            if self.max_length and text and len(text) > self.max_length:
                logger.debug(f"Truncating {file_path.name}: too long ({len(text)} chars)")
                text = text[:self.max_length]
            
            return text
            
        except DataError:
            # Re-raise DataError exceptions (like missing dependencies)
            raise
        except Exception as e:
            logger.error(f"Error loading {file_path}: {str(e)}")
            return None
    
    def _load_txt(self, file_path: Path) -> str:
        """Load text from a TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise DataError(f"Could not decode {file_path} with any supported encoding")
    
    def _load_pdf(self, file_path: Path) -> str:
        """Load text from a PDF file."""
        if not HAS_PDF:
            raise DataError("PyPDF2 not installed. Install with: pip install PyPDF2")
        
        try:
            text_content = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(page_text)
                    except Exception as e:
                        logger.warning(f"Error extracting page {page_num} from {file_path}: {e}")
                        continue
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise DataError(f"Error reading PDF {file_path}: {str(e)}")
    
    def _load_docx(self, file_path: Path) -> str:
        """Load text from a DOCX file."""
        if not HAS_DOCX:
            raise DataError("python-docx not installed. Install with: pip install python-docx")
        
        try:
            doc = Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise DataError(f"Error reading DOCX {file_path}: {str(e)}")
    
    def _load_pptx(self, file_path: Path) -> str:
        """Load text from a PPTX file."""
        if not HAS_PPTX:
            raise DataError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            raise DataError(f"Error reading PPTX {file_path}: {str(e)}")
    
    def _load_markdown(self, file_path: Path) -> str:
        """Load text from a Markdown file."""
        # For now, just load as plain text
        # Could add markdown parsing later if needed
        return self._load_txt(file_path)
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text."""
        if not text:
            return text
        
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # Multiple newlines to double
        text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces/tabs to single space
        text = re.sub(r'^\s+|\s+$', '', text, flags=re.MULTILINE)  # Strip line whitespace
        
        # Remove common artifacts
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)  # Control chars
        text = re.sub(r'�', '', text)  # Replacement character
        
        return text.strip()
    
    def load_directory(self, 
                      directory: Union[str, Path],
                      recursive: bool = True,
                      pattern: Optional[str] = None) -> Iterator[tuple[Path, str]]:
        """
        Load all supported files from a directory.
        
        Args:
            directory: Directory to scan
            recursive: Whether to scan subdirectories
            pattern: Optional filename pattern to match
            
        Yields:
            tuple[Path, str]: File path and extracted text content
        """
        directory = Path(directory)
        
        if not directory.exists():
            raise DataError(f"Directory not found: {directory}")
        
        if not directory.is_dir():
            raise DataError(f"Path is not a directory: {directory}")
        
        # Find files
        if recursive:
            files = directory.rglob("*")
        else:
            files = directory.glob("*")
        
        # Filter by extension and pattern
        for file_path in files:
            if not file_path.is_file():
                continue
            
            if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
                continue
            
            if pattern and not re.search(pattern, file_path.name):
                continue
            
            try:
                text = self.load_file(file_path)
                if text:  # Only yield if text was successfully extracted
                    yield file_path, text
            except Exception as e:
                logger.error(f"Error loading {file_path}: {e}")
                continue
    
    def get_file_metadata(self, file_path: Union[str, Path], text: Optional[str] = None) -> DocumentMetadata:
        """
        Get metadata for a file.
        
        Args:
            file_path: Path to the file
            text: Optional pre-loaded text content
            
        Returns:
            DocumentMetadata: File metadata
        """
        file_path = Path(file_path)
        
        if text is None:
            text = self.load_file(file_path) or ""
        
        return DocumentMetadata(
            filename=file_path.name,
            file_path=str(file_path.absolute()),
            file_size=file_path.stat().st_size if file_path.exists() else 0,
            file_type=file_path.suffix.lower(),
            char_count=len(text),
            word_count=len(text.split()) if text else 0,
            line_count=text.count('\n') + 1 if text else 0
        )
    
    def get_supported_extensions(self) -> set[str]:
        """Get set of supported file extensions."""
        return self.SUPPORTED_EXTENSIONS.copy()
    
    def is_supported(self, file_path: Union[str, Path]) -> bool:
        """Check if a file format is supported."""
        return Path(file_path).suffix.lower() in self.SUPPORTED_EXTENSIONS